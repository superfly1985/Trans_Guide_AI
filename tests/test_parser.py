# -*- coding: utf-8 -*-
"""测试文件解析器"""

from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt
import sys
sys.path.insert(0, 'd:\\01.AwesomeProject\\52.Trans_Guide_AI')

# 创建测试 Excel 文件
print("创建测试 Excel 文件...")
wb = Workbook()
ws = wb.active
ws.title = 'Sheet1'

# 添加一些测试数据
ws['A1'] = 'Hello World'
ws['B1'] = 'Test Document'
ws['A2'] = 'This is a sample text for translation'
ws['B2'] = 'Formula result'
ws['C2'] = '=A1&B1'  # 公式

wb.save('tests/test_excel.xlsx')
print('测试 Excel 文件已创建: tests/test_excel.xlsx')

# 创建测试 PPT 文件
print("\n创建测试 PPT 文件...")
prs = Presentation()
slide_layout = prs.slide_layouts[6]  # 空白布局

# 添加第一张幻灯片
slide1 = prs.slides.add_slide(slide_layout)
title_box = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "Welcome to the Presentation"

content_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.text = "This is a sample text for translation testing."

# 添加第二张幻灯片
slide2 = prs.slides.add_slide(slide_layout)
title_box2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
title_frame2 = title_box2.text_frame
title_frame2.text = "Second Slide"

content_box2 = slide2.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
content_frame2 = content_box2.text_frame
content_frame2.text = "More content to translate here."

prs.save('tests/test_pptx.pptx')
print('测试 PPT 文件已创建: tests/test_pptx.pptx')

# 测试解析
print("\n" + "="*50)
print("测试 Excel 解析:")
print("="*50)
from modules.file_parser import parse_file
blocks, info = parse_file('tests/test_excel.xlsx')
print(f'解析成功! 文本块数量: {len(blocks)}')
print(f'文件信息: {info}')
for block in blocks:
    print(f'  块 {block["index"]}: Sheet={block["sheet"]}, Row={block["row"]}, Col={block["col"]}, Text="{block["text"]}"')

print("\n" + "="*50)
print("测试 PPT 解析:")
print("="*50)
blocks, info = parse_file('tests/test_pptx.pptx')
print(f'解析成功! 文本块数量: {len(blocks)}')
print(f'文件信息: {info}')
for block in blocks:
    print(f'  块 {block["index"]}: Slide={block["slide"]}, Shape={block["shape_id"]}, Text="{block["text"][:50]}..."')

print("\n" + "="*50)
print("测试 .ppt 旧格式拒绝:")
print("="*50)
try:
    from modules.file_parser import parse_ppt_file
    parse_ppt_file('tests/test.ppt')
except ValueError as e:
    print(f'正确拒绝: {e}')

print("\n所有测试完成!")
