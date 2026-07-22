# -*- coding: utf-8 -*-
"""
文件解析模块
支持Word、Excel、PDF、TXT、CSV格式
"""

import os
import csv
from typing import List, Dict, Tuple, Optional


def parse_word_file(file_path: str) -> Tuple[List[Dict], Dict]:
    """
    解析Word文件
    
    Args:
        file_path: 文件路径
        
    Returns:
        texts: 文本块列表，每个元素包含:
            - text: 文本内容
            - type: 类型 (paragraph/table/header/footer)
            - style: 段落样式名
            - index: 在文档中的索引
            - format: 格式信息 (字体、字号、加粗等)
        format_info: 文档整体格式信息
        
    Raises:
        FileNotFoundError: 文件不存在
        ValueError: 文件格式错误或损坏
    """
    try:
        from docx import Document
    except ImportError:
        raise ImportError("请安装 python-docx: pip install python-docx")
    
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在: {file_path}")
    
    try:
        doc = Document(file_path)
    except Exception as e:
        raise ValueError(f"无法解析Word文件: {e}")
    
    texts = []
    index = 0
    seen_texts = set()
    
    # 解析段落
    for para_idx, para in enumerate(doc.paragraphs):
        text = para.text.strip()
        if text and text not in seen_texts:
            seen_texts.add(text)
            # 提取格式信息
            format_info = _extract_para_format(para)
            
            texts.append({
                "text": text,
                "type": "paragraph",
                "style": para.style.name if para.style else "Normal",
                "index": index,
                "para_index": para_idx,
                "format": format_info
            })
            index += 1
    
    # 解析表格
    for table_idx, table in enumerate(doc.tables):
        for row_idx, row in enumerate(table.rows):
            for cell_idx, cell in enumerate(row.cells):
                text = cell.text.strip()
                if text and text not in seen_texts:
                    seen_texts.add(text)
                    texts.append({
                        "text": text,
                        "type": "table_cell",
                        "table_index": table_idx,
                        "row": row_idx,
                        "col": cell_idx,
                        "index": index,
                        "format": {}
                    })
                    index += 1
    
    # 解析页眉
    for section_idx, section in enumerate(doc.sections):
        if section.header:
            for para in section.header.paragraphs:
                text = para.text.strip()
                if text and text not in seen_texts:
                    seen_texts.add(text)
                    texts.append({
                        "text": text,
                        "type": "header",
                        "section_index": section_idx,
                        "index": index,
                        "format": {}
                    })
                    index += 1
    
    # 解析页脚
    for section_idx, section in enumerate(doc.sections):
        if section.footer:
            for para in section.footer.paragraphs:
                text = para.text.strip()
                if text and text not in seen_texts:
                    seen_texts.add(text)
                    texts.append({
                        "text": text,
                        "type": "footer",
                        "section_index": section_idx,
                        "index": index,
                        "format": {}
                    })
                    index += 1
    
    doc_info = {
        "paragraph_count": len(doc.paragraphs),
        "table_count": len(doc.tables),
        "section_count": len(doc.sections),
        "text_block_count": len(texts)
    }
    
    return texts, doc_info


def _extract_para_format(para) -> Dict:
    """
    提取段落格式信息
    
    Args:
        para: docx段落对象
        
    Returns:
        格式信息字典
    """
    format_info = {
        "alignment": str(para.alignment) if para.alignment else None,
        "line_spacing": para.paragraph_format.line_spacing,
        "space_before": para.paragraph_format.space_before,
        "space_after": para.paragraph_format.space_after,
    }
    
    # 提取第一个run的格式（如果有）
    if para.runs:
        run = para.runs[0]
        format_info.update({
            "bold": run.bold,
            "italic": run.italic,
            "underline": run.underline,
            "font_name": run.font.name if run.font else None,
            "font_size": run.font.size.pt if run.font and run.font.size else None,
        })
    
    return format_info


def parse_excel_file(file_path: str) -> Tuple[List[Dict], Dict]:
    """
    解析Excel文件（支持 .xlsx, .xlsm）
    注意: .xls 格式不支持，请转换为 .xlsx
    
    Args:
        file_path: 文件路径
        
    Returns:
        texts: 单元格文本列表，每个元素包含:
            - text: 文本内容
            - sheet: 工作表名
            - row: 行号 (0-based)
            - col: 列号 (0-based)
            - index: 在文档中的索引
            - is_formula: 是否是公式
            - format: 格式信息
        format_info: 工作簿整体格式信息
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在: {file_path}")
    
    ext = os.path.splitext(file_path)[1].lower()
    
    # 拒绝 .xls 格式
    if ext == '.xls':
        raise ValueError("不支持 .xls 格式，请先将文件另存为 .xlsx 格式后再上传")
    
    # 根据扩展名选择解析器
    if ext in ['.xlsx', '.xlsm']:
        return _parse_excel_openpyxl(file_path)
    else:
        raise ValueError(f"不支持的Excel格式: {ext}")


MAX_STANDARD_EXCEL_CELLS = 500_000


def _parse_excel_openpyxl(file_path: str) -> Tuple[List[Dict], Dict]:
    """使用 openpyxl 解析 .xlsx 和 .xlsm 文件"""
    try:
        from openpyxl import load_workbook
        from .special_parsers import find_china_sheet_name, parse_china_sheet
    except ImportError:
        raise ImportError("请安装 openpyxl: pip install openpyxl")
    
    try:
        wb = load_workbook(file_path, data_only=False)  # data_only=False 以检测公式
    except Exception as e:
        raise ValueError(f"无法解析Excel文件: {e}")
    
    try:
        china_sheet_name = find_china_sheet_name(wb.sheetnames)
        if china_sheet_name:
            return parse_china_sheet(file_path, china_sheet_name)
        
        # 标准解析逻辑
        texts = []
        index = 0
        total_cells = 0
        formula_count = 0
        skipped_large_sheets = []
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_cell_count = ws.max_row * ws.max_column
            if sheet_cell_count > MAX_STANDARD_EXCEL_CELLS:
                skipped_large_sheets.append(sheet_name)
                cells = sorted(ws._cells.values(), key=lambda cell: (cell.row, cell.column))
            else:
                cells = (cell for row in ws.iter_rows() for cell in row)
            
            for cell in cells:
                if cell.value is not None:
                    # 检测是否是公式
                    is_formula = cell.data_type == 'f' or (isinstance(cell.value, str) and cell.value.startswith('='))
                    
                    # 如果是公式，跳过不翻译
                    if is_formula:
                        formula_count += 1
                        continue
                    
                    cell_text = str(cell.value).strip()
                    if cell_text:
                        # 提取格式信息
                        format_info = _extract_cell_format(cell)
                        
                        texts.append({
                            "text": cell_text,
                            "sheet": sheet_name,
                            "row": cell.row - 1,
                            "col": cell.column - 1,
                            "index": index,
                            "is_formula": False,
                            "format": format_info
                        })
                        index += 1
                        total_cells += 1
        
        wb_info = {
            "sheet_count": len(wb.sheetnames),
            "sheet_names": wb.sheetnames,
            "cell_count": total_cells,
            "formula_count": formula_count
        }
        if skipped_large_sheets:
            wb_info["skipped_large_sheets"] = skipped_large_sheets
        
        return texts, wb_info
    finally:
        wb.close()


def _parse_china_sheet(wb, file_path: str) -> Tuple[List[Dict], Dict]:
    """
    解析特殊结构的 "中国" sheet
    只提取 "English" 列的内容（读取公式计算后的值），翻译后将结果填入 "Landessprache / Local Language" 列
    """
    from openpyxl import load_workbook
    
    # 使用 data_only=True 重新打开文件，获取公式计算后的值
    try:
        wb_data = load_workbook(file_path, data_only=True)
        ws = wb_data["中国"]
    except Exception as e:
        raise ValueError(f"无法读取公式计算值: {e}")
    
    texts = []
    index = 0
    
    # 找到标题行，确定 "English" 和 "Landessprache / Local Language" 列的位置
    english_col = None
    local_lang_col = None
    header_row = 0
    
    # 遍历前5行查找标题
    for row_idx in range(min(5, ws.max_row)):
        row = list(ws.iter_rows(min_row=row_idx+1, max_row=row_idx+1))[0]
        for col_idx, cell in enumerate(row):
            if cell.value:
                header_text = str(cell.value).strip()
                if header_text == "English":
                    english_col = col_idx
                elif header_text == "Landessprache / Local Language":
                    local_lang_col = col_idx
        
        # 如果找到了两个列，停止搜索
        if english_col is not None and local_lang_col is not None:
            header_row = row_idx
            break
    
    # 如果找不到标题列，使用默认的 B 列和 C 列
    if english_col is None:
        english_col = 1  # B 列（索引从0开始，所以1是B列）
        logger.warning(f'[ChinaSheet] 未找到 "English" 列，使用默认 B 列 (列索引: {english_col})')

    if local_lang_col is None:
        local_lang_col = 2  # C 列（索引从0开始，所以2是C列）
        logger.warning(f'[ChinaSheet] 未找到 "Landessprache / Local Language" 列，使用默认 C 列 (列索引: {local_lang_col})')
    
    # 提取 "English" 列的数据（从标题行下一行开始）
    # 读取公式计算后的值，而不是公式本身
    for row_idx in range(header_row + 1, ws.max_row):
        english_cell = ws.cell(row=row_idx + 1, column=english_col + 1)
        
        # 使用 data_only=True 读取的值已经是计算后的结果
        if english_cell.value is not None:
            cell_text = str(english_cell.value).strip()
            if cell_text and cell_text != "English":  # 跳过标题本身
                # 对于这种特殊表格，不需要保留格式信息
                texts.append({
                    "text": cell_text,
                    "sheet": "中国",
                    "row": row_idx,
                    "col": local_lang_col,  # 目标列是 "Landessprache / Local Language"
                    "source_col": english_col,  # 记录源列
                    "index": index,
                    "is_formula": False,  # 已经是计算后的值
                    "format": {},  # 简化格式信息
                    "is_china_sheet": True  # 标记为特殊结构
                })
                index += 1
    
    wb_info = {
        "sheet_count": len(wb.sheetnames),
        "sheet_names": wb.sheetnames,
        "cell_count": len(texts),
        "formula_count": 0,
        "special_structure": "china_sheet",
        "english_col": english_col,
        "local_lang_col": local_lang_col,
        "header_row": header_row
    }
    
    return texts, wb_info


def _parse_excel_xlrd(file_path: str) -> Tuple[List[Dict], Dict]:
    """使用 xlrd 解析 .xls 文件"""
    try:
        import xlrd
    except ImportError:
        raise ImportError("请安装 xlrd: pip install xlrd")
    
    try:
        wb = xlrd.open_workbook(file_path)
    except Exception as e:
        raise ValueError(f"无法解析Excel文件: {e}")
    
    texts = []
    index = 0
    total_cells = 0
    
    for sheet_idx in range(wb.nsheets):
        ws = wb.sheet_by_index(sheet_idx)
        sheet_name = ws.name
        
        for row_idx in range(ws.nrows):
            for col_idx in range(ws.ncols):
                cell = ws.cell(row_idx, col_idx)
                
                # xlrd: 0=empty, 1=text, 2=number, 3=date, 4=boolean, 5=error, 6=blank
                if cell.ctype == xlrd.XL_CELL_EMPTY or cell.ctype == xlrd.XL_CELL_BLANK:
                    continue
                
                # 跳过公式（xlrd 读取的是计算结果，但可以通过其他方式检测）
                # xlrd 在读取 .xls 时，公式已经计算为值，无法直接检测
                # 我们假设非空值都需要翻译
                
                cell_text = str(cell.value).strip()
                if cell_text:
                    texts.append({
                        "text": cell_text,
                        "sheet": sheet_name,
                        "row": row_idx,
                        "col": col_idx,
                        "index": index,
                        "is_formula": False,  # xlrd 无法直接检测公式
                        "format": {}
                    })
                    index += 1
                    total_cells += 1
    
    wb_info = {
        "sheet_count": wb.nsheets,
        "sheet_names": wb.sheet_names(),
        "cell_count": total_cells
    }
    
    return texts, wb_info


def _extract_cell_format(cell) -> Dict:
    """
    提取单元格格式信息
    
    Args:
        cell: openpyxl单元格对象
        
    Returns:
        格式信息字典
    """
    format_info = {}
    
    if cell.font:
        format_info.update({
            "font_name": cell.font.name,
            "font_size": cell.font.size,
            "bold": cell.font.bold,
            "italic": cell.font.italic,
            "underline": cell.font.underline,
        })
    
    if cell.fill and cell.fill.fgColor:
        # 将 RGB 对象转换为字符串
        rgb_value = cell.fill.fgColor.rgb
        if rgb_value is not None:
            format_info["fill_color"] = str(rgb_value)
    
    if cell.alignment:
        format_info.update({
            "horizontal": cell.alignment.horizontal,
            "vertical": cell.alignment.vertical,
            "wrap_text": cell.alignment.wrap_text,
        })
    
    if cell.border:
        format_info["has_border"] = any([
            cell.border.left.style,
            cell.border.right.style,
            cell.border.top.style,
            cell.border.bottom.style
        ])
    
    return format_info


def parse_pdf_file(file_path: str) -> Tuple[List[Dict], Dict]:
    """
    解析PDF文件（仅文本提取）
    
    Args:
        file_path: 文件路径
        
    Returns:
        texts: 文本块列表，每个元素包含:
            - text: 文本内容
            - page: 页码 (0-based)
            - index: 在文档中的索引
            - bbox: 文本块位置 (x0, y0, x1, y1)
        format_info: 文档页数等信息
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise ImportError("请安装 PyMuPDF: pip install PyMuPDF")
    
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在: {file_path}")
    
    try:
        doc = fitz.open(file_path)
    except Exception as e:
        raise ValueError(f"无法解析PDF文件: {e}")
    
    texts = []
    index = 0
    
    for page_num in range(len(doc)):
        page = doc[page_num]
        
        # 提取文本块
        text_blocks = page.get_text("blocks")
        
        for block in text_blocks:
            # block格式: (x0, y0, x1, y1, text, block_no, block_type)
            x0, y0, x1, y1, text, block_no, block_type = block
            
            text = text.strip()
            if text:
                texts.append({
                    "text": text,
                    "page": page_num,
                    "index": index,
                    "bbox": (x0, y0, x1, y1),
                    "block_no": block_no,
                    "block_type": block_type
                })
                index += 1
    
    doc_info = {
        "page_count": len(doc),
        "text_block_count": len(texts)
    }
    
    doc.close()
    return texts, doc_info


def parse_txt_file(file_path: str) -> Tuple[List[Dict], Dict]:
    """
    解析TXT文件
    
    Args:
        file_path: 文件路径
        
    Returns:
        texts: 文本行列表
        format_info: 文件编码等信息
    """
    texts = []
    encodings = ["utf-8", "gbk", "gb2312", "utf-16"]
    
    for encoding in encodings:
        try:
            with open(file_path, "r", encoding=encoding) as f:
                lines = f.readlines()
            for i, line in enumerate(lines):
                texts.append({
                    "text": line.rstrip("\n\r"),
                    "line_no": i + 1,
                    "type": "line"
                })
            return texts, {"encoding": encoding, "line_count": len(lines)}
        except UnicodeDecodeError:
            continue
    
    raise ValueError(f"无法识别文件编码: {file_path}")


def parse_csv_file(file_path: str) -> Tuple[List[Dict], Dict]:
    """
    解析CSV文件
    
    Args:
        file_path: 文件路径
        
    Returns:
        texts: 单元格文本列表
        format_info: 行列数等信息
    """
    texts = []
    encodings = ["utf-8", "gbk", "gb2312"]
    
    for encoding in encodings:
        try:
            with open(file_path, "r", encoding=encoding, newline="") as f:
                reader = csv.reader(f)
                for row_idx, row in enumerate(reader):
                    for col_idx, cell in enumerate(row):
                        texts.append({
                            "text": cell,
                            "row": row_idx,
                            "col": col_idx,
                            "type": "cell"
                        })
            return texts, {"encoding": encoding, "rows": row_idx + 1}
        except UnicodeDecodeError:
            continue
    
    raise ValueError(f"无法识别文件编码: {file_path}")


def parse_doc_file(file_path: str) -> Tuple[List[Dict], Dict]:
    """
    .doc 格式不支持，请转换为 .docx
    """
    raise ValueError(
        "不支持 .doc 格式，请先将文件另存为 .docx 格式后再上传"
    )


def parse_pptx_file(file_path: str) -> Tuple[List[Dict], Dict]:
    """
    解析PowerPoint文件 (.pptx)
    
    Args:
        file_path: 文件路径
        
    Returns:
        texts: 文本块列表，每个元素包含:
            - text: 文本内容
            - slide: 幻灯片编号 (0-based)
            - shape_id: 形状ID
            - shape_type: 形状类型
            - index: 在文档中的索引
        format_info: 演示文稿整体信息
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("请安装 python-pptx: pip install python-pptx")
    
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在: {file_path}")
    
    try:
        prs = Presentation(file_path)
    except Exception as e:
        raise ValueError(f"无法解析PPT文件: {e}")
    
    texts = []
    index = 0
    total_shapes = 0
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            # 检查形状是否有文本框
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                if text:
                    texts.append({
                        "text": text,
                        "slide": slide_idx,
                        "shape_id": shape.shape_id,
                        "shape_type": str(shape.shape_type),
                        "index": index,
                        "format": {
                            "left": shape.left,
                            "top": shape.top,
                            "width": shape.width,
                            "height": shape.height
                        }
                    })
                    index += 1
                    total_shapes += 1
            
            # 处理表格
            if shape.has_table:
                table = shape.table
                for row_idx, row in enumerate(table.rows):
                    for cell_idx, cell in enumerate(row.cells):
                        if cell.text.strip():
                            texts.append({
                                "text": cell.text.strip(),
                                "slide": slide_idx,
                                "shape_id": shape.shape_id,
                                "shape_type": "table_cell",
                                "table_row": row_idx,
                                "table_col": cell_idx,
                                "index": index,
                                "format": {}
                            })
                            index += 1
                            total_shapes += 1
    
    ppt_info = {
        "slide_count": len(prs.slides),
        "text_block_count": len(texts),
        "shape_count": total_shapes
    }
    
    return texts, ppt_info


def parse_ppt_file(file_path: str) -> Tuple[List[Dict], Dict]:
    """
    处理旧版 .ppt 文件 - 拒绝并提示用户转换
    """
    raise ValueError(
        "不支持 .ppt 格式，请先将文件另存为 .pptx 格式后再上传"
    )


def parse_file(file_path: str) -> Tuple[List[Dict], Dict]:
    """
    自动识别文件类型并解析
    
    Args:
        file_path: 文件路径
        
    Returns:
        texts: 文本内容列表
        format_info: 格式信息
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在: {file_path}")
    
    ext = os.path.splitext(file_path)[1].lower()
    
    parsers = {
        ".docx": parse_word_file,
        ".doc": parse_doc_file,
        ".xlsx": parse_excel_file,
        ".xlsm": parse_excel_file,
        ".pptx": parse_pptx_file,
        ".ppt": parse_ppt_file,
        ".pdf": parse_pdf_file,
        ".txt": parse_txt_file,
        ".csv": parse_csv_file
    }
    
    if ext in parsers:
        return parsers[ext](file_path)
    else:
        raise ValueError(f"不支持的文件格式: {ext}")


def get_file_type(file_path: str) -> str:
    """
    获取文件类型
    
    Args:
        file_path: 文件路径
        
    Returns:
        文件类型标识
    """
    ext = os.path.splitext(file_path)[1].lower()
    type_map = {
        ".docx": "word",
        ".doc": "word",
        ".xlsx": "excel",
        ".xls": "excel",
        ".xlsm": "excel",
        ".pptx": "ppt",
        ".ppt": "ppt",
        ".pdf": "pdf",
        ".txt": "txt",
        ".csv": "csv"
    }
    return type_map.get(ext, "unknown")
