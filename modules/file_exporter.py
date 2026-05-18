# -*- coding: utf-8 -*-
"""
文件导出模块
将译文写回Word/Excel文件
"""

import os
from typing import List, Tuple, Dict


def export_word(
    original_path: str,
    translations: List[Tuple[int, str]],
    output_path: str,
    mode: str = "bilingual"
) -> bool:
    """
    导出翻译后的Word文件
    
    Args:
        original_path: 原文件路径
        translations: 翻译结果列表 [(位置索引, 译文), ...]
        位置索引对应 parse_word_file 返回的 index 字段
        output_path: 输出路径
        mode: "bilingual"（双语）| "target_only"（仅译文）
        
    Returns:
        是否成功
    """
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor
    except ImportError:
        raise ImportError("请安装 python-docx: pip install python-docx")
    
    if not os.path.exists(original_path):
        raise FileNotFoundError(f"原文件不存在: {original_path}")
    
    try:
        # 打开原文档
        doc = Document(original_path)
        
        # 构建翻译映射
        trans_map = {idx: text for idx, text in translations}
        
        # 处理段落
        para_idx = 0
        for para in doc.paragraphs:
            if para.text.strip() and para_idx in trans_map:
                translation = trans_map[para_idx]
                
                if mode == "bilingual":
                    # 双语模式：原文 + 换行 + 译文
                    # 保留原文，在后面添加译文
                    para.add_run("\n")
                    run = para.add_run(translation)
                    # 设置译文格式（红色以示区分）
                    run.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
                else:
                    # 仅译文模式：替换原文
                    para.clear()
                    para.add_run(translation)
                
                para_idx += 1
        
        # 处理表格
        table_idx_offset = para_idx
        for table_idx, table in enumerate(doc.tables):
            for row_idx, row in enumerate(table.rows):
                for cell_idx, cell in enumerate(row.cells):
                    current_idx = table_idx_offset + table_idx * 1000 + row_idx * 100 + cell_idx
                    if current_idx in trans_map:
                        translation = trans_map[current_idx]
                        
                        if mode == "bilingual":
                            cell.paragraphs[0].add_run("\n")
                            run = cell.paragraphs[0].add_run(translation)
                            run.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
                        else:
                            cell.paragraphs[0].clear()
                            cell.paragraphs[0].add_run(translation)
        
        # 保存文档
        doc.save(output_path)
        return True
        
    except Exception as e:
        print(f"导出Word失败: {e}")
        return False


def export_word_simple(
    original_path: str,
    translation_dict: Dict[int, str],
    output_path: str,
    mode: str = "bilingual"
) -> bool:
    """
    简化版Word导出（使用字典映射）
    
    Args:
        original_path: 原文件路径
        translation_dict: 翻译字典 {index: 译文}
        output_path: 输出路径
        mode: "bilingual" | "target_only"
        
    Returns:
        是否成功
    """
    try:
        from docx import Document
        from docx.shared import RGBColor
    except ImportError:
        raise ImportError("请安装 python-docx: pip install python-docx")
    
    if not os.path.exists(original_path):
        raise FileNotFoundError(f"原文件不存在: {original_path}")
    
    try:
        doc = Document(original_path)
        
        # 收集所有文本块及其位置
        text_blocks = []
        index = 0
        seen_texts = set()
        
        # 段落
        for para in doc.paragraphs:
            text = para.text.strip()
            if text and text not in seen_texts:
                seen_texts.add(text)
                text_blocks.append(("paragraph", para, index))
                index += 1
        
        # 表格
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text = cell.text.strip()
                    if text and text not in seen_texts:
                        seen_texts.add(text)
                        text_blocks.append(("cell", cell, index))
                        index += 1
        
        print(f"[导出调试] 找到 {len(text_blocks)} 个文本块")
        print(f"[导出调试] 翻译字典包含 {len(translation_dict)} 个翻译")
        print(f"[导出调试] 文本块索引: {[idx for _, _, idx in text_blocks[:10]]}...")
        print(f"[导出调试] 翻译字典键: {list(translation_dict.keys())[:10]}...")
        
        def _copy_format(source_runs, target_run):
            if source_runs:
                src = source_runs[0]
                if src.font.name:
                    target_run.font.name = src.font.name
                if src.font.size:
                    target_run.font.size = src.font.size
                target_run.bold = src.bold
                target_run.italic = src.italic
                target_run.underline = src.underline
                if src.font.color and src.font.color.rgb:
                    target_run.font.color.rgb = src.font.color.rgb
        
        def _dedup_translation(original_text, translated_text):
            if not original_text or not translated_text:
                return translated_text
            orig_stripped = original_text.strip()
            trans_stripped = translated_text.strip()
            if trans_stripped == orig_stripped:
                return orig_stripped
            if trans_stripped.startswith(orig_stripped):
                remaining = trans_stripped[len(orig_stripped):].strip()
                if remaining:
                    return remaining
            return translated_text

        def _text_already_contains(element, chinese_text):
            existing = element.text
            lines = [l.strip() for l in existing.split('\n') if l.strip()]
            target = chinese_text.strip()
            return any(target in line for line in lines)
        
        # 应用翻译
        applied_count = 0
        for block_type, block, idx in text_blocks:
            if idx in translation_dict:
                translation = translation_dict[idx]
                
                if block_type == "paragraph":
                    para = block
                    original_text = para.text.strip()
                    translation = _dedup_translation(original_text, translation)

                    if _text_already_contains(para, translation):
                        print(f"[导出调试] 段落 {idx}: 译文已存在，跳过")
                        continue

                    if mode == "bilingual":
                        para.add_run("\n")
                        run = para.add_run(translation)
                        _copy_format(para.runs[:1], run)
                        run.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
                    else:
                        runs_before = list(para.runs)
                        para.clear()
                        run = para.add_run(translation)
                        _copy_format(runs_before, run)
                    applied_count += 1
                
                elif block_type == "cell":
                    cell = block
                    cell_para = cell.paragraphs[0]
                    original_text = cell_para.text.strip()
                    translation = _dedup_translation(original_text, translation)

                    if _text_already_contains(cell_para, translation):
                        print(f"[导出调试] 单元格 {idx}: 译文已存在，跳过")
                        continue

                    if mode == "bilingual":
                        cell_para.add_run("\n")
                        run = cell_para.add_run(translation)
                        _copy_format(cell_para.runs[:1], run)
                        run.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
                    else:
                        runs_before = list(cell_para.runs)
                        cell_para.clear()
                        run = cell_para.add_run(translation)
                        _copy_format(runs_before, run)
                    applied_count += 1
            else:
                print(f"[导出调试] 索引 {idx} 没有对应的翻译")
        
        print(f"[导出调试] 成功应用 {applied_count} 个翻译")
        doc.save(output_path)
        return True
        
    except Exception as e:
        print(f"导出Word失败: {e}")
        import traceback
        print(traceback.format_exc())
        return False
        return False


def export_excel(
    original_path: str,
    translations: List[Tuple[str, int, int, str]],
    output_path: str,
    mode: str = "bilingual"
) -> bool:
    """
    导出翻译后的Excel文件
    
    Args:
        original_path: 原文件路径
        translations: 翻译结果列表 [(sheet名, 行, 列, 译文), ...]
        行、列对应 parse_excel_file 返回的 row、col 字段
        output_path: 输出路径
        mode: "bilingual" | "target_only"
        
    Returns:
        是否成功
    """
    try:
        from openpyxl import load_workbook
        from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
    except ImportError:
        raise ImportError("请安装 openpyxl: pip install openpyxl")
    
    if not os.path.exists(original_path):
        raise FileNotFoundError(f"原文件不存在: {original_path}")
    
    try:
        # 打开原工作簿
        wb = load_workbook(original_path)
        
        # 构建翻译映射 {(sheet, row, col): 译文}
        trans_map = {}
        for sheet_name, row, col, translation in translations:
            trans_map[(sheet_name, row, col)] = translation
        
        # 应用翻译
        for sheet_name, row, col, translation in translations:
            if sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                cell = ws.cell(row=row + 1, column=col + 1)  # openpyxl使用1-based索引
                
                if mode == "bilingual":
                    # 双语模式：原文 + 换行 + 译文
                    original_value = str(cell.value) if cell.value else ""
                    cell.value = f"{original_value}\n{translation}"
                    # 设置译文为红色
                    cell.font = Font(color="FF0000", italic=True)
                else:
                    # 仅译文模式：替换原文
                    cell.value = translation
        
        # 保存工作簿
        wb.save(output_path)
        return True
        
    except Exception as e:
        print(f"导出Excel失败: {e}")
        return False


def export_excel_simple(
    original_path: str,
    translation_dict: Dict[Tuple[str, int, int], str],
    output_path: str,
    mode: str = "bilingual",
    is_china_sheet: bool = False
) -> bool:
    """
    简化版Excel导出（使用字典映射）
    
    Args:
        original_path: 原文件路径
        translation_dict: 翻译字典 {(sheet, row, col): 译文}
        output_path: 输出路径
        mode: "bilingual" | "target_only"
        is_china_sheet: 是否是特殊结构的 "中国" sheet
        
    Returns:
        是否成功
    """
    try:
        from openpyxl import load_workbook
        from openpyxl.styles import Font
    except ImportError:
        raise ImportError("请安装 openpyxl: pip install openpyxl")
    
    if not os.path.exists(original_path):
        raise FileNotFoundError(f"原文件不存在: {original_path}")
    
    try:
        wb = load_workbook(original_path)
        
        # 特殊处理 "中国" sheet 结构
        if is_china_sheet and "中国" in wb.sheetnames:
            ws = wb["中国"]
            for (sheet_name, row, col), translation in translation_dict.items():
                if sheet_name == "中国":
                    # 直接写入 "Landessprache / Local Language" 列
                    cell = ws.cell(row=row + 1, column=col + 1)
                    cell.value = translation
                    cell.font = Font(color="FF0000", italic=True)
        else:
            # 标准导出逻辑
            for (sheet_name, row, col), translation in translation_dict.items():
                if sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    cell = ws.cell(row=row + 1, column=col + 1)
                    
                    if mode == "bilingual":
                        original_value = str(cell.value) if cell.value else ""
                        cell.value = f"{original_value}\n{translation}"
                        cell.font = Font(color="FF0000", italic=True)
                    else:
                        cell.value = translation
        
        wb.save(output_path)
        return True
        
    except Exception as e:
        print(f"导出Excel失败: {e}")
        return False


def export_txt(
    translations: List[Dict],
    output_path: str
) -> bool:
    """
    导出翻译后的TXT文件
    
    Args:
        translations: 翻译结果列表
        output_path: 输出路径
        
    Returns:
        是否成功
    """
    try:
        with open(output_path, "w", encoding="utf-8") as f:
            for item in translations:
                f.write(item.get("translation", "") + "\n")
        return True
    except Exception as e:
        print(f"导出TXT失败: {e}")
        return False


def export_csv(
    translations: List[Dict],
    output_path: str
) -> bool:
    """
    导出翻译后的CSV文件
    
    Args:
        translations: 翻译结果列表
        output_path: 输出路径
        
    Returns:
        是否成功
    """
    import csv
    
    try:
        with open(output_path, "w", encoding="utf-8", newline="") as f:
            writer = csv.writer(f)
            for item in translations:
                writer.writerow([
                    item.get("original", ""),
                    item.get("translation", "")
                ])
        return True
    except Exception as e:
        print(f"导出CSV失败: {e}")
        return False


def export_pptx(
    original_path: str,
    translation_dict: Dict[int, str],
    output_path: str,
    mode: str = "bilingual"
) -> bool:
    """
    导出翻译后的PPT文件
    
    Args:
        original_path: 原文件路径
        translation_dict: 翻译字典 {index: 译文}
        output_path: 输出路径
        mode: "bilingual" | "target_only"
        
    Returns:
        是否成功
    """
    try:
        from pptx import Presentation
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        raise ImportError("请安装 python-pptx: pip install python-pptx")
    
    if not os.path.exists(original_path):
        raise FileNotFoundError(f"原文件不存在: {original_path}")
    
    try:
        prs = Presentation(original_path)
        
        text_blocks = []
        index = 0
        
        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_blocks.append(("shape", shape, index))
                    index += 1
                
                if shape.has_table:
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        for cell_idx, cell in enumerate(row.cells):
                            if cell.text.strip():
                                text_blocks.append(("table_cell", cell, index))
                                index += 1
        
        for block_type, block, idx in text_blocks:
            if idx in translation_dict:
                translation = translation_dict[idx]
                
                if block_type == "shape":
                    shape = block
                    if mode == "bilingual":
                        if hasattr(shape, "text_frame"):
                            last_para = shape.text_frame.paragraphs[-1]
                            new_para = shape.text_frame.add_paragraph()
                            new_run = new_para.add_run()
                            new_run.text = translation
                            new_run.font.color.rgb = RGBColor(255, 0, 0)
                    else:
                        if hasattr(shape, "text_frame"):
                            for para in shape.text_frame.paragraphs:
                                for run in para.runs:
                                    run.text = ""
                            shape.text_frame.paragraphs[0].runs[0].text = translation
                
                elif block_type == "table_cell":
                    cell = block
                    if mode == "bilingual":
                        original_text = cell.text
                        cell.text = f"{original_text}\n{translation}"
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.color.rgb = RGBColor(255, 0, 0)
                    else:
                        cell.text = translation
        
        prs.save(output_path)
        return True
        
    except Exception as e:
        print(f"导出PPT失败: {e}")
        return False


def get_output_path(input_path: str, suffix: str = "_中文") -> str:
    """
    生成输出文件路径
    
    Args:
        input_path: 输入文件路径
        suffix: 文件名后缀
        
    Returns:
        输出文件路径
    """
    dir_name = os.path.dirname(input_path)
    base_name = os.path.basename(input_path)
    name, ext = os.path.splitext(base_name)
    output_name = f"{name}{suffix}{ext}"
    return os.path.join(dir_name, output_name)
